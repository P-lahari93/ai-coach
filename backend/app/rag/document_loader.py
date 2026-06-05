# FILE: app/rag/document_loader.py
"""
DocumentLoader — load documents from various sources.

Supported formats:
  - PDF (pypdf)
  - DOCX (python-docx)
  - PPTX (python-pptx)
  - TXT (plain text)
  - Markdown (.md)
  - URL (httpx + beautifulsoup4)

All methods return the extracted text content as a string.
"""
from __future__ import annotations

import mimetypes
from pathlib import Path
from typing import Tuple

import httpx
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

from app.core.exceptions import IngestionError


class DocumentLoader:
    """Load and extract text from various document formats."""

    def __init__(self) -> None:
        """Initialize the document loader with timeout configuration."""
        self._http_timeout = 30.0

    async def load_pdf(self, path: str) -> str:
        """
        Extract text from a PDF file using pypdf.

        Args:
            path: absolute file path to the PDF

        Returns:
            extracted text concatenated from all pages

        Raises:
            IngestionError: when PDF cannot be read or is corrupted
        """
        try:
            reader = PdfReader(path)
            pages_text: list[str] = []
            
            for page in reader.pages:
                text = page.extract_text()
                if text and text.strip():
                    pages_text.append(text)
            
            if not pages_text:
                raise IngestionError(f"No text content extracted from PDF: {path}")
            
            return "\n\n".join(pages_text)
        except Exception as exc:
            raise IngestionError(
                f"Failed to load PDF from {path}: {exc}"
            ) from exc

    async def load_docx(self, path: str) -> str:
        """
        Extract text from a DOCX file using python-docx.

        Args:
            path: absolute file path to the DOCX

        Returns:
            extracted text concatenated from all paragraphs

        Raises:
            IngestionError: when DOCX cannot be read or is corrupted
        """
        try:
            doc = DocxDocument(path)
            paragraphs_text: list[str] = []
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs_text.append(text)
            
            if not paragraphs_text:
                raise IngestionError(f"No text content extracted from DOCX: {path}")
            
            return "\n\n".join(paragraphs_text)
        except Exception as exc:
            raise IngestionError(
                f"Failed to load DOCX from {path}: {exc}"
            ) from exc

    async def load_pptx(self, path: str) -> str:
        """
        Extract text from a PPTX file using python-pptx.

        Extracts text from all shapes in all slides, maintaining
        slide order. Slide titles and body text are both included.

        Args:
            path: absolute file path to the PPTX

        Returns:
            extracted text concatenated from all slides

        Raises:
            IngestionError: when PPTX cannot be read or is corrupted
        """
        try:
            prs = Presentation(path)
            slides_text: list[str] = []
            
            for slide in prs.slides:
                slide_parts: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text.strip())
                
                if slide_parts:
                    slides_text.append("\n".join(slide_parts))
            
            if not slides_text:
                raise IngestionError(f"No text content extracted from PPTX: {path}")
            
            return "\n\n".join(slides_text)
        except Exception as exc:
            raise IngestionError(
                f"Failed to load PPTX from {path}: {exc}"
            ) from exc

    async def load_text(self, path: str) -> str:
        """
        Load plain text file.

        Args:
            path: absolute file path to the text file

        Returns:
            file contents as string

        Raises:
            IngestionError: when file cannot be read
        """
        try:
            with open(path, "r", encoding="utf-8") as f:
                content = f.read()
            
            if not content.strip():
                raise IngestionError(f"Empty text file: {path}")
            
            return content
        except Exception as exc:
            raise IngestionError(
                f"Failed to load text file from {path}: {exc}"
            ) from exc

    async def load_markdown(self, path: str) -> str:
        """
        Load Markdown file as plain text.

        For MVP, Markdown is treated as plain text. Future enhancement:
        parse markdown structure and extract headings for better chunking.

        Args:
            path: absolute file path to the markdown file

        Returns:
            file contents as string

        Raises:
            IngestionError: when file cannot be read
        """
        return await self.load_text(path)

    async def load_from_url(self, url: str) -> Tuple[str, str]:
        """
        Fetch and extract main content from a URL.

        Uses httpx for HTTP fetching and BeautifulSoup for HTML parsing.
        Extracts visible text from <p>, <article>, <section>, <div> tags,
        removing scripts, styles, and navigation elements.

        Args:
            url: the URL to fetch

        Returns:
            tuple of (content: str, mime_type: str)

        Raises:
            IngestionError: when URL cannot be fetched or parsed
        """
        try:
            async with httpx.AsyncClient(timeout=self._http_timeout) as client:
                response = await client.get(url, follow_redirects=True)
                response.raise_for_status()
                
                mime_type = response.headers.get("content-type", "text/html").split(";")[0]
                
                if "html" not in mime_type.lower():
                    raise IngestionError(
                        f"URL returned non-HTML content: {mime_type}"
                    )
                
                html = response.text
                soup = BeautifulSoup(html, "html.parser")
                
                # Remove script and style elements
                for element in soup(["script", "style", "nav", "header", "footer"]):
                    element.decompose()
                
                # Extract text from main content areas
                content_tags = soup.find_all(
                    ["p", "article", "section", "div", "h1", "h2", "h3", "h4", "h5", "h6"]
                )
                
                text_parts: list[str] = []
                for tag in content_tags:
                    text = tag.get_text(separator=" ", strip=True)
                    if text and len(text) > 20:  # filter out navigation/UI text
                        text_parts.append(text)
                
                if not text_parts:
                    raise IngestionError(f"No meaningful content extracted from URL: {url}")
                
                content = "\n\n".join(text_parts)
                return content, mime_type
                
        except httpx.HTTPError as exc:
            raise IngestionError(
                f"Failed to fetch URL {url}: {exc}"
            ) from exc
        except Exception as exc:
            raise IngestionError(
                f"Failed to parse content from {url}: {exc}"
            ) from exc

    async def load(self, path: str, mime_type: str | None = None) -> str:
        """
        Router method to load a document using the appropriate loader.

        Args:
            path: absolute file path
            mime_type: optional MIME type hint; if None, inferred from extension

        Returns:
            extracted text content

        Raises:
            IngestionError: when format is unsupported or loading fails
        """
        if mime_type is None:
            mime_type, _ = mimetypes.guess_type(path)
        
        # Normalize mime_type
        if mime_type:
            mime_type = mime_type.lower()
        
        # Map MIME types to loaders
        if mime_type == "application/pdf" or path.endswith(".pdf"):
            return await self.load_pdf(path)
        elif mime_type in (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
        ) or path.endswith(".docx"):
            return await self.load_docx(path)
        elif mime_type in (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ) or path.endswith(".pptx"):
            return await self.load_pptx(path)
        elif mime_type == "text/markdown" or path.endswith(".md"):
            return await self.load_markdown(path)
        elif mime_type and mime_type.startswith("text/") or path.endswith(".txt"):
            return await self.load_text(path)
        else:
            raise IngestionError(
                f"Unsupported document format: {mime_type or 'unknown'} (path: {path})"
            )
